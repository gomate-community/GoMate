#!/usr/bin/env python
# -*- coding:utf-8 _*-
"""
@author:quincy qiang
@license: Apache Licence
@file: views.py
@time: 2024/06/13
@contact: yanqiangmiffy@gamil.com
@software: PyCharm
@description: coding..
"""
from typing import List

from pydantic import BaseModel, Field
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from fastapi import APIRouter
import magic
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
from typing import List

parse_router = APIRouter()

@parse_router.post("/parse/", response_model=None, summary="文件解析")
async def parser(file: UploadFile = File(...), description: str = None):
    try:
        # 读取文件内容
        content = await file.read()

        # 检测文件类型
        mime = magic.Magic(mime=True)
        file_type = mime.from_buffer(content)

        if file_type == "text/plain" or file.filename.endswith('.md'):
            # 处理txt或markdown文件
            text = content.decode("utf-8")
            # 在这里处理文本内容，例如保存到数据库
            with open(f"/path/to/save/{file.filename}", "w", encoding="utf-8") as f:
                f.write(text)

        elif file_type == "application/pdf":
            # 处理PDF文件
            with pdfplumber.open(file.file) as pdf:
                text = "\n".join(page.extract_text() for page in pdf.pages)
                # 在这里处理文本内容，例如保存到数据库
                with open(f"/path/to/save/{file.filename}.txt", "w", encoding="utf-8") as f:
                    f.write(text)

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # 处理docx文件
            doc = docx.Document(file.file)
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            # 在这里处理文本内容，例如保存到数据库
            with open(f"/path/to/save/{file.filename}.txt", "w", encoding="utf-8") as f:
                f.write(text)

        elif file_type == "text/html":
            # 处理html文件
            text = content.decode("utf-8")
            # 在这里处理HTML内容，例如解析和保存到数据库
            with open(f"/path/to/save/{file.filename}", "w", encoding="utf-8") as f:
                f.write(text)

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # 处理ppt文件
            prs = Presentation(file.file)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            # 在这里处理文本内容，例如保存到数据库
            with open(f"/path/to/save/{file.filename}.txt", "w", encoding="utf-8") as f:
                f.write(text)

        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # 处理excel文件
            wb = openpyxl.load_workbook(file.file)
            sheet = wb.active
            text = "\n".join(" ".join(cell.value for cell in row) for row in sheet.iter_rows())
            # 在这里处理文本内容，例如保存到数据库
            with open(f"/path/to/save/{file.filename}.txt", "w", encoding="utf-8") as f:
                f.write(text)

        else:
            raise HTTPException(status_code=400, detail="不支持的文件类型")

        # 返回成功响应
        return JSONResponse(content={"filename": file.filename, "description": description, "file_type": file_type}, status_code=200)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件上传失败: {str(e)}")

